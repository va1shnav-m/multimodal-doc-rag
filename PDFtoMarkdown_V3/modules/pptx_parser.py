import io
from pathlib import Path
from typing import Dict, Any, List
from PIL import Image

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None


def _extract_table(shape) -> str:
    """Convert a PPTX table shape into a GitHub-flavored Markdown table."""
    table = shape.table
    rows = []
    for r_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if r_idx == 0:
            separator = "| " + " | ".join(["---"] * len(cells)) + " |"
            rows.append(separator)
    return "\n".join(rows)


def parse_pptx(
    pptx_path: str | Path,
    assets_dir: str | Path,
    document_prefix: str = "slide",
) -> Dict[str, Any]:
    """
    Parse a PowerPoint (.pptx) file into structured Markdown.
    Extracts text, bullet points, tables, speaker notes, and embedded images.

    Parameters
    ----------
    pptx_path : str | Path
        Path to the .pptx file.
    assets_dir : str | Path
        Directory to save extracted images.
    document_prefix : str
        Prefix for saved image filenames.

    Returns
    -------
    dict
        {
            "markdown": str,
            "image_count": int,
            "slide_count": int,
            "images": List[Path]
        }
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is not installed. "
            "Please install it using: pip install python-pptx"
        )

    pptx_path = Path(pptx_path)
    assets_dir = Path(assets_dir)
    assets_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    slides_md = []
    extracted_images = []
    image_counter = 0

    slides_md.append(f"# Presentation: {pptx_path.name}\n")

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        slide_title = None

        # 1. Check title placeholder
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip().replace("\n", " ")

        if slide_title:
            slide_lines.append(f"## Slide {slide_num}: {slide_title}\n")
        else:
            slide_lines.append(f"## Slide {slide_num}\n")

        # 2. Iterate through shapes
        for shape in slide.shapes:
            # Skip title shape since it's already rendered
            if shape == slide.shapes.title:
                continue

            # Tables
            if shape.has_table:
                table_md = _extract_table(shape)
                if table_md:
                    slide_lines.append(table_md)
                    slide_lines.append("")
                continue

            # Images / Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                try:
                    image_bytes = shape.image.blob
                    image_ext = shape.image.ext or "png"
                    image_counter += 1
                    img_filename = f"{document_prefix}_{slide_num:03d}_{image_counter:03d}.{image_ext}"
                    img_path = assets_dir / img_filename

                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    extracted_images.append(img_path)
                    slide_lines.append(f"![Slide {slide_num} Image](assets/{img_filename})\n")
                except Exception as e:
                    print(f"Warning: Failed to extract image on slide {slide_num}: {e}")
                continue

            # Text boxes / Shapes with text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    level = getattr(paragraph, "level", 0)
                    indent = "  " * level

                    # Render as bullet if indented or shape has multiple paragraphs
                    if level > 0 or len(shape.text_frame.paragraphs) > 1:
                        slide_lines.append(f"{indent}- {text}")
                    else:
                        slide_lines.append(f"{text}\n")

        # 3. Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_lines.append("\n> **Speaker Notes:**")
                for note_line in notes_text.splitlines():
                    if note_line.strip():
                        slide_lines.append(f"> {note_line.strip()}")
                slide_lines.append("")

        slides_md.append("\n".join(slide_lines).strip())
        slides_md.append("\n\n---\n\n")

    full_markdown = "".join(slides_md).strip()

    return {
        "markdown": full_markdown,
        "image_count": len(extracted_images),
        "slide_count": len(prs.slides),
        "images": extracted_images,
    }
